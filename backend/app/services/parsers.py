"""
AtlasLM source parsers  -  DOCX and CSV support (Patch 002).

Each parser returns the same shape the pipeline already consumes:
    List[{"page_number": int, "content": str}]

DOCX: paragraphs + tables, grouped into synthetic "pages" of ~3000 chars
      (DOCX has no real pages without rendering; synthetic pages keep the
      citation drawer meaningful: "Section N" instead of a fake page).
CSV:  header-aware row serialization, grouped in blocks of rows so each
      chunk keeps column context ("col: value" pairs survive chunking).
"""

import csv
import io
import logging
from typing import List, Dict, Any

logger = logging.getLogger("atlaslm.parsers")

# Synthetic page size for paginating non-paged formats
SECTION_CHAR_TARGET = 3000
# CSV rows serialized per section
CSV_ROWS_PER_SECTION = 50
CSV_MAX_FIELD_LEN = 500


def extract_text_from_docx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """Extract paragraphs and tables from a .docx file in document order."""
    try:
        from docx import Document as DocxDocument  # python-docx
    except ImportError as exc:
        raise ValueError(
            "DOCX support is not installed on the server (python-docx missing)."
        ) from exc

    logger.info("Parsing DOCX: %s (%d bytes)", filename, len(file_bytes))
    try:
        doc = DocxDocument(io.BytesIO(file_bytes))
    except Exception as exc:
        raise ValueError(
            f"Could not read '{filename}'. The file may be corrupted or not a valid Word document."
        ) from exc

    blocks: List[str] = []

    # Paragraphs and tables in document body order
    from docx.table import Table
    from docx.text.paragraph import Paragraph
    from docx.oxml.ns import qn

    body = doc.element.body
    for child in body.iterchildren():
        if child.tag == qn("w:p"):
            text = Paragraph(child, doc).text.strip()
            if text:
                blocks.append(text)
        elif child.tag == qn("w:tbl"):
            table = Table(child, doc)
            rows_out = []
            for row in table.rows:
                cells = [c.text.strip().replace("\n", " ") for c in row.cells]
                if any(cells):
                    rows_out.append(" | ".join(cells))
            if rows_out:
                blocks.append("[TABLE]\n" + "\n".join(rows_out) + "\n[/TABLE]")

    if not blocks:
        raise ValueError(
            f"No extractable text found in '{filename}'. The document appears to be empty."
        )

    # Group blocks into synthetic sections (~SECTION_CHAR_TARGET chars)
    pages: List[Dict[str, Any]] = []
    current: List[str] = []
    current_len = 0
    section = 1
    for block in blocks:
        current.append(block)
        current_len += len(block)
        if current_len >= SECTION_CHAR_TARGET:
            pages.append({"page_number": section, "content": "\n\n".join(current)})
            section += 1
            current, current_len = [], 0
    if current:
        pages.append({"page_number": section, "content": "\n\n".join(current)})

    logger.info("DOCX parsed: %s -> %d sections", filename, len(pages))
    return pages


def extract_text_from_csv(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Serialize CSV rows with header context so every chunk remains
    self-describing for retrieval ("Region: EMEA; Revenue: 1.2M; ...").
    """
    logger.info("Parsing CSV: %s (%d bytes)", filename, len(file_bytes))

    text = None
    for encoding in ("utf-8-sig", "utf-8", "latin-1"):
        try:
            text = file_bytes.decode(encoding)
            break
        except UnicodeDecodeError:
            continue
    if text is None:
        raise ValueError(f"Could not decode '{filename}' as text.")

    # Sniff delimiter (comma/semicolon/tab) with sane fallback
    sample = text[:8192]
    try:
        dialect = csv.Sniffer().sniff(sample, delimiters=",;\t|")
    except csv.Error:
        dialect = csv.excel

    reader = csv.reader(io.StringIO(text), dialect)
    rows = [r for r in reader if any(field.strip() for field in r)]
    if not rows:
        raise ValueError(f"'{filename}' contains no data rows.")

    header = [h.strip() or f"column_{i+1}" for i, h in enumerate(rows[0])]
    data_rows = rows[1:] if len(rows) > 1 else []

    if not data_rows:
        # Header-only file: still index the schema
        return [{
            "page_number": 1,
            "content": f"CSV file '{filename}' with columns: {', '.join(header)} (no data rows).",
        }]

    pages: List[Dict[str, Any]] = []
    section = 1
    for start in range(0, len(data_rows), CSV_ROWS_PER_SECTION):
        block_rows = data_rows[start:start + CSV_ROWS_PER_SECTION]
        lines = [
            f"CSV: {filename} | Columns: {', '.join(header)} | "
            f"Rows {start + 1}-{start + len(block_rows)} of {len(data_rows)}"
        ]
        for offset, row in enumerate(block_rows):
            pairs = []
            for i, value in enumerate(row):
                value = value.strip()
                if not value:
                    continue
                if len(value) > CSV_MAX_FIELD_LEN:
                    value = value[:CSV_MAX_FIELD_LEN] + "..."
                col = header[i] if i < len(header) else f"column_{i+1}"
                pairs.append(f"{col}: {value}")
            if pairs:
                lines.append("; ".join(pairs))
        pages.append({"page_number": section, "content": "\n".join(lines)})
        section += 1

    logger.info("CSV parsed: %s -> %d rows in %d sections", filename, len(data_rows), len(pages))
    return pages


# ====================================================================== #
# XLSX (Patch 003)
# ====================================================================== #

XLSX_ROWS_PER_SECTION = 50
XLSX_MAX_FIELD_LEN = 500
XLSX_MAX_CELLS = 500_000  # hard guard against pathological spreadsheets


def extract_text_from_xlsx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extract an XLSX workbook into synthetic sections.

    - Each sheet is processed independently.
    - The first non-empty row of a sheet is treated as the header.
    - Rows are rendered as "Header: value" pairs (same convention as CSV)
      so retrieval keeps column context in every chunk.
    - Sections of XLSX_ROWS_PER_SECTION rows; page_number increments
      globally across sheets so citations stay unique per document.
    """
    import io
    try:
        from openpyxl import load_workbook
    except ImportError as e:
        raise ValueError("XLSX support is not available on this server.") from e

    logger.info("Starting XLSX extraction: %s (%d bytes)", filename, len(file_bytes))

    try:
        wb = load_workbook(io.BytesIO(file_bytes), read_only=True, data_only=True)
    except Exception as e:
        logger.error("XLSX parser failure for %s: %s", filename, e)
        raise ValueError(f"Failed to parse XLSX '{filename}': file may be corrupt or password-protected.")

    pages: List[Dict[str, Any]] = []
    section = 1
    total_cells = 0

    def _cell(v: Any) -> str:
        if v is None:
            return ""
        s = str(v).strip()
        if len(s) > XLSX_MAX_FIELD_LEN:
            s = s[:XLSX_MAX_FIELD_LEN] + "..."
        return s

    for sheet in wb.worksheets:
        rows_iter = sheet.iter_rows(values_only=True)
        header: List[str] = []
        data_rows: List[List[str]] = []

        for raw in rows_iter:
            values = [_cell(v) for v in raw]
            if not any(values):
                continue
            total_cells += len(values)
            if total_cells > XLSX_MAX_CELLS:
                logger.warning("XLSX %s truncated at %d cells.", filename, XLSX_MAX_CELLS)
                break
            if not header:
                header = [v if v else f"column_{i+1}" for i, v in enumerate(values)]
                continue
            data_rows.append(values)

        if not header:
            continue

        if not data_rows:
            pages.append({
                "page_number": section,
                "content": f"XLSX: {filename} | Sheet: {sheet.title} | Columns: {', '.join(header)} (no data rows).",
            })
            section += 1
            continue

        for start in range(0, len(data_rows), XLSX_ROWS_PER_SECTION):
            block = data_rows[start:start + XLSX_ROWS_PER_SECTION]
            lines = [
                f"XLSX: {filename} | Sheet: {sheet.title} | Columns: {', '.join(header)} | "
                f"Rows {start + 1}-{start + len(block)} of {len(data_rows)}"
            ]
            for offset, row in enumerate(block):
                pairs = []
                for i, value in enumerate(row):
                    if not value:
                        continue
                    col = header[i] if i < len(header) else f"column_{i+1}"
                    pairs.append(f"{col}: {value}")
                if pairs:
                    lines.append("; ".join(pairs))
            pages.append({"page_number": section, "content": "\n".join(lines)})
            section += 1

    wb.close()

    if not pages:
        raise ValueError(f"No extractable data found in XLSX '{filename}'.")

    logger.info("XLSX parsed: %s -> %d sections across %d sheet(s)", filename, len(pages), len(wb.sheetnames))
    return pages


# ====================================================================== #
# PPTX (Patch 003)
# ====================================================================== #

def extract_text_from_pptx(file_bytes: bytes, filename: str) -> List[Dict[str, Any]]:
    """
    Extract a PPTX deck into per-slide sections.

    Captures, in order: slide title, body text frames, table contents,
    and speaker notes. page_number == slide number, so citations render
    naturally as 'Slide N' in the UI.
    """
    import io
    try:
        from pptx import Presentation
    except ImportError as e:
        raise ValueError("PPTX support is not available on this server.") from e

    logger.info("Starting PPTX extraction: %s (%d bytes)", filename, len(file_bytes))

    try:
        prs = Presentation(io.BytesIO(file_bytes))
    except Exception as e:
        logger.error("PPTX parser failure for %s: %s", filename, e)
        raise ValueError(f"Failed to parse PPTX '{filename}': file may be corrupt or password-protected.")

    pages: List[Dict[str, Any]] = []

    for slide_idx, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []

        title = ""
        if slide.shapes.title is not None and slide.shapes.title.has_text_frame:
            title = slide.shapes.title.text_frame.text.strip()
        parts.append(f"Slide {slide_idx}" + (f": {title}" if title else ""))

        for shape in slide.shapes:
            if shape == slide.shapes.title:
                continue
            if shape.has_text_frame:
                text = "\n".join(
                    p.text.strip() for p in shape.text_frame.paragraphs if p.text.strip()
                )
                if text:
                    parts.append(text)
            if getattr(shape, "has_table", False) and shape.has_table:
                for row in shape.table.rows:
                    cells = [c.text.strip() for c in row.cells]
                    if any(cells):
                        parts.append(" | ".join(cells))

        if slide.has_notes_slide:
            notes = slide.notes_slide.notes_text_frame.text.strip()
            if notes:
                parts.append(f"Speaker notes: {notes}")

        content = "\n".join(parts).strip()
        # Always emit the slide section, even if sparse, to keep slide
        # numbering aligned with the source deck.
        pages.append({"page_number": slide_idx, "content": content})

    non_empty = [p for p in pages if len(p["content"].split("\n")) > 1 or ":" in p["content"]]
    if not non_empty:
        raise ValueError(
            f"No extractable text found in PPTX '{filename}'. "
            "The deck may contain only images."
        )

    logger.info("PPTX parsed: %s -> %d slides", filename, len(pages))
    return pages
