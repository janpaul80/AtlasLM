# backend/app/services/ingest/pptx_loader.py
"""PPTX loader -> per-slide text + speaker notes. Requires python-pptx."""
from __future__ import annotations
from typing import List
from pptx import Presentation
from .base import ExtractedBlock, block


def load_pptx(path: str) -> List[ExtractedBlock]:
    prs = Presentation(path)
    blocks: List[ExtractedBlock] = []
    for i, slide in enumerate(prs.slides, start=1):
        parts: List[str] = []
        for shape in slide.shapes:
            if shape.has_text_frame and shape.text_frame.text.strip():
                parts.append(shape.text_frame.text.strip())
        notes = ""
        if slide.has_notes_slide and slide.notes_slide.notes_text_frame:
            notes = slide.notes_slide.notes_text_frame.text.strip()
        body = "\n".join(parts)
        if notes:
            body += f"\n[Speaker notes] {notes}"
        if body.strip():
            blocks.append(block(body, page=i))  # 'page' == slide number
    return blocks
